#!/usr/bin/env python3
"""
Charterwell Investor/Sales Pitch Deck Generator
Generates a brand-compliant .pptx deck following the KEN-49 spec.

Colors: Deep Navy (#1B2A4A), Teal (#0D9488), Warm White (#FAFAF8), Charcoal (#1F2937)
Typography: Inter family (Semibold headlines, Regular body)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors
DEEP_NAVY = RGBColor(0x1B, 0x2A, 0x4A)
TEAL = RGBColor(0x0D, 0x94, 0x88)
WARM_WHITE = RGBColor(0xFA, 0xFA, 0xF8)
CHARCOAL = RGBColor(0x1F, 0x29, 0x37)
LIGHT_TEAL = RGBColor(0xCC, 0xFB, 0xF1)  # teal tint for backgrounds
SLATE = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Slide dimensions (widescreen 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


def add_bg(slide, color):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 font_color=CHARCOAL, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Inter", line_spacing=1.2):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    return txBox


def add_multi_text(slide, left, top, width, height, lines, font_size=16,
                   font_color=CHARCOAL, bold=False, alignment=PP_ALIGN.LEFT,
                   font_name="Inter", bullet=False, line_spacing=Pt(24)):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if bullet and line.strip():
            p.text = line
            p.level = 0
        else:
            p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = line_spacing
    return txBox


def add_divider(slide, left, top, width, color=TEAL):
    """Add a thin horizontal line."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DEEP_NAVY)

# Accent line
add_rect(slide, Inches(1.5), Inches(2.2), Inches(1.5), Pt(4), TEAL)

# Company name
add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.2),
             "CHARTERWELL", font_size=54, font_color=WHITE, bold=True)

# Tagline
add_text_box(slide, Inches(1.5), Inches(3.7), Inches(10), Inches(0.8),
             "The AI-Native Claims Intelligence Workspace",
             font_size=28, font_color=TEAL, bold=False)

# Sub-tagline
add_text_box(slide, Inches(1.5), Inches(4.8), Inches(8), Inches(0.6),
             "Every document understood. Every decision informed.",
             font_size=18, font_color=WARM_WHITE, bold=False)

# Bottom bar
add_rect(slide, Inches(0), Inches(7.0), SLIDE_WIDTH, Inches(0.5), TEAL)
add_text_box(slide, Inches(1.5), Inches(7.05), Inches(6), Inches(0.4),
             "Confidential  |  March 2026", font_size=12, font_color=DEEP_NAVY, bold=False)


# ============================================================
# SLIDE 2: THE PROBLEM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WARM_WHITE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.5),
             "THE PROBLEM", font_size=14, font_color=TEAL, bold=True)
add_divider(slide, Inches(1), Inches(1.0), Inches(2))

add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1),
             "Claims professionals are drowning in documents.",
             font_size=36, font_color=DEEP_NAVY, bold=True)

# Stats row - three columns
stats = [
    ("10\u201315", "documents per claim,\nprocessed manually"),
    ("17", "disconnected data sources\nper carrier"),
    ("15\u201320 min", "per FNOL just to validate\nwhat should be obvious"),
]
for i, (num, desc) in enumerate(stats):
    x = Inches(1 + i * 3.8)
    # Stat box
    add_rect(slide, x, Inches(2.8), Inches(3.3), Inches(2.2), WHITE, SLATE)
    add_text_box(slide, x + Inches(0.3), Inches(3.0), Inches(2.8), Inches(0.8),
                 num, font_size=42, font_color=TEAL, bold=True, alignment=PP_ALIGN.LEFT)
    add_text_box(slide, x + Inches(0.3), Inches(3.9), Inches(2.8), Inches(1.0),
                 desc, font_size=14, font_color=CHARCOAL, bold=False)

# Bottom insight
add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(1.2),
             "Point solutions for each step create a fragmented stack.\n"
             "Legacy platforms bolt AI onto decades-old architecture.\n"
             "The result: 80% of carriers deploy AI. Only 7% scale it.",
             font_size=18, font_color=CHARCOAL, bold=False)


# ============================================================
# SLIDE 3: THE SHIFT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WARM_WHITE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.5),
             "THE SHIFT", font_size=14, font_color=TEAL, bold=True)
add_divider(slide, Inches(1), Inches(1.0), Inches(2))

add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1),
             "From point solutions to unified workspaces.",
             font_size=36, font_color=DEEP_NAVY, bold=True)

add_multi_text(slide, Inches(1), Inches(2.5), Inches(5.5), Inches(3.5), [
    '"Claims Intelligence" is the emerging category.',
    "",
    "Parallel precedents:",
    "\u2022  Legora created \u201cThe AI Workspace for Lawyers\u201d \u2192 $5.55B valuation",
    "\u2022  Veeva created \u201cLife Sciences Cloud\u201d \u2192 $50B+ market cap",
    "\u2022  Figma created \u201cDesign Workspace\u201d \u2192 $20B (Adobe acquisition)",
    "\u2022  Bloomberg defined the \u201cTerminal\u201d for traders \u2192 category dominance",
    "",
    "The pattern: professionals who live in documents",
    "and data get a purpose-built workspace.",
    "",
    "Insurance claims is next.",
], font_size=16, font_color=CHARCOAL)

# Right side callout
add_rect(slide, Inches(7.5), Inches(2.5), Inches(5), Inches(3.5), DEEP_NAVY)
add_multi_text(slide, Inches(7.8), Inches(2.8), Inches(4.5), Inches(3.0), [
    "Industry signals:",
    "",
    "\u2022  88% of carriers want AI",
    "\u2022  Only 7% have scaled it",
    "\u2022  SAS: \u201c2026 is the year AI becomes",
    "    the operating system for insurance\u201d",
    "\u2022  A Fortune 500 insurer is phasing",
    "    out policy admin for AI copilots",
    "",
    "The category window is open now.",
], font_size=15, font_color=WHITE)


# ============================================================
# SLIDE 4: THE INSIGHT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DEEP_NAVY)

add_text_box(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.5),
             "THE INSIGHT", font_size=14, font_color=TEAL, bold=True)
add_divider(slide, Inches(1), Inches(1.0), Inches(2), TEAL)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
             "The document layer is\nthe foundation.",
             font_size=44, font_color=WHITE, bold=True)

add_multi_text(slide, Inches(1), Inches(3.5), Inches(5.5), Inches(3.5), [
    "Every claims decision starts with",
    "understanding documents.",
    "",
    "\u2022  FNOL reports",
    "\u2022  Police reports",
    "\u2022  Medical records",
    "\u2022  Repair estimates",
    "\u2022  Policy declarations",
    "\u2022  Correspondence",
    "",
    "No one owns this layer today.",
], font_size=18, font_color=WARM_WHITE)

# Right callout
add_rect(slide, Inches(7.5), Inches(3.5), Inches(5), Inches(2.5), TEAL)
add_multi_text(slide, Inches(7.8), Inches(3.7), Inches(4.5), Inches(2.2), [
    "The document layer connects to",
    "every downstream decision:",
    "",
    "Coverage determination",
    "Adjuster routing",
    "Fraud detection",
    "Compliance checking",
    "Subrogation",
    "Settlement",
], font_size=15, font_color=WHITE)


# ============================================================
# SLIDE 5: THE PRODUCT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WARM_WHITE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.5),
             "THE PRODUCT", font_size=14, font_color=TEAL, bold=True)
add_divider(slide, Inches(1), Inches(1.0), Inches(2))

add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1),
             "One workspace replaces your claims AI stack.",
             font_size=36, font_color=DEEP_NAVY, bold=True)

# Three pillars
pillars = [
    ("Document Intelligence\nEngine", "AI that doesn\u2019t just extract data \u2014\nit understands what documents mean\nfor your claim. 50+ P&C document\ntypes. Context-aware extraction."),
    ("Playbook System", "Carrier-specific rules encoded as\nliving playbooks. Adjudication logic,\ncompliance requirements, escalation\ntriggers. Your institutional knowledge,\nsystemized."),
    ("Unified Workspace", "One environment where claims\nprofessionals think, decide, and act.\nNo more alt-tabbing between 12\ntools. Every document, decision, and\nworkflow \u2014 connected."),
]
for i, (title, desc) in enumerate(pillars):
    x = Inches(1 + i * 3.8)
    add_rect(slide, x, Inches(2.5), Inches(3.4), Inches(4.2), WHITE, SLATE)
    # Accent bar at top
    add_rect(slide, x, Inches(2.5), Inches(3.4), Pt(5), TEAL)
    add_text_box(slide, x + Inches(0.3), Inches(2.8), Inches(2.9), Inches(1.0),
                 title, font_size=20, font_color=DEEP_NAVY, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(4.0), Inches(2.9), Inches(2.5),
                 desc, font_size=14, font_color=CHARCOAL)


# ============================================================
# SLIDE 6: HOW IT WORKS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WARM_WHITE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.5),
             "HOW IT WORKS", font_size=14, font_color=TEAL, bold=True)
add_divider(slide, Inches(1), Inches(1.0), Inches(2))

add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
             "From document to decision in minutes, not hours.",
             font_size=32, font_color=DEEP_NAVY, bold=True)

# Workflow steps
steps = [
    ("1", "INGEST", "Documents flow in from\nemail, uploads, or\ncore system APIs."),
    ("2", "UNDERSTAND", "AI reads, classifies, and\nextracts \u2014 understanding\ncontext, not just text."),
    ("3", "ROUTE", "Claims triaged and\nassigned to the right\nadjuster automatically."),
    ("4", "CHECK", "Compliance verified\nagainst 50-state rules\nand carrier playbooks."),
    ("5", "ACT", "Adjusters work in a\nunified workspace with\nfull claim context."),
]
for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.5 + i * 2.5)
    # Circle with number
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.7), Inches(2.6), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = TEAL
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)

    add_text_box(slide, x + Inches(0.2), Inches(3.5), Inches(2.2), Inches(0.5),
                 title, font_size=16, font_color=DEEP_NAVY, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(4.1), Inches(2.3), Inches(1.5),
                 desc, font_size=13, font_color=CHARCOAL, alignment=PP_ALIGN.CENTER)

    # Arrow between steps
    if i < 4:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(2.2), Inches(2.8), Inches(0.4), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = SLATE
        arrow.line.fill.background()

# Bottom callout
add_rect(slide, Inches(1), Inches(5.8), Inches(11.3), Inches(1.2), DEEP_NAVY)
add_text_box(slide, Inches(1.5), Inches(5.95), Inches(10.5), Inches(0.9),
             "The entire pipeline is AI-native \u2014 not AI bolted on. Every step learns from the last.\n"
             "Playbooks encode your carrier-specific rules so the workspace gets smarter with every claim.",
             font_size=15, font_color=WHITE, alignment=PP_ALIGN.LEFT)


# ============================================================
# SLIDE 7: KEY METRICS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WARM_WHITE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.5),
             "KEY METRICS", font_size=14, font_color=TEAL, bold=True)
add_divider(slide, Inches(1), Inches(1.0), Inches(2))

add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
             "Measurable impact from day one.",
             font_size=36, font_color=DEEP_NAVY, bold=True)

metrics = [
    ("85%", "Processing time\nreduction", "Beating the Hyperscience\nbenchmark for top-10 insurers"),
    ("< 3 min", "FNOL triage time", "Down from 15\u201320 minutes\nper first notice of loss"),
    ("40%", "Claims cycle time\nreduction", "End-to-end from FNOL\nto resolution"),
    ("95%+", "Extraction accuracy", "Across 50+ P&C\ndocument types"),
]
for i, (num, title, sub) in enumerate(metrics):
    x = Inches(0.8 + i * 3.1)
    add_rect(slide, x, Inches(2.5), Inches(2.8), Inches(3.5), WHITE, SLATE)
    add_text_box(slide, x + Inches(0.2), Inches(2.7), Inches(2.4), Inches(1.0),
                 num, font_size=48, font_color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(3.8), Inches(2.4), Inches(0.8),
                 title, font_size=18, font_color=DEEP_NAVY, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(4.8), Inches(2.4), Inches(1.0),
                 sub, font_size=13, font_color=SLATE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.6),
             "Pilot framework: 60\u201390 days, one claim type, one geography. Measurable ROI before enterprise commitment.",
             font_size=16, font_color=CHARCOAL, alignment=PP_ALIGN.LEFT)


# ============================================================
# SLIDE 8: COMPETITIVE LANDSCAPE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WARM_WHITE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.5),
             "COMPETITIVE LANDSCAPE", font_size=14, font_color=TEAL, bold=True)
add_divider(slide, Inches(1), Inches(1.0), Inches(2))

add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
             "Deep intelligence meets unified workspace.",
             font_size=32, font_color=DEEP_NAVY, bold=True)

# 2x2 matrix background
matrix_left = Inches(1.5)
matrix_top = Inches(2.5)
matrix_w = Inches(7)
matrix_h = Inches(4.5)

# Quadrant backgrounds
half_w = Inches(3.5)
half_h = Inches(2.25)
add_rect(slide, matrix_left, matrix_top, half_w, half_h, RGBColor(0xF0, 0xFD, 0xFA))  # top-left
add_rect(slide, matrix_left + half_w, matrix_top, half_w, half_h, RGBColor(0xCC, 0xFB, 0xF1))  # top-right (Charterwell zone)
add_rect(slide, matrix_left, matrix_top + half_h, half_w, half_h, RGBColor(0xF8, 0xFA, 0xFC))  # bottom-left
add_rect(slide, matrix_left + half_w, matrix_top + half_h, half_w, half_h, RGBColor(0xF0, 0xFD, 0xFA))  # bottom-right

# Axis labels
add_text_box(slide, Inches(1.0), Inches(4.5), Inches(0.5), Inches(1.0),
             "POINT\nSOLUTION", font_size=10, font_color=SLATE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(8.2), Inches(4.5), Inches(0.8), Inches(1.0),
             "UNIFIED\nWORKSPACE", font_size=10, font_color=SLATE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(4.5), Inches(2.1), Inches(1.5), Inches(0.4),
             "DEEP AI", font_size=10, font_color=SLATE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(4.5), Inches(7.0), Inches(1.5), Inches(0.4),
             "SHALLOW AI", font_size=10, font_color=SLATE, bold=True, alignment=PP_ALIGN.CENTER)

# Competitor dots
competitors = [
    ("Tractable", Inches(2.2), Inches(3.0), SLATE),
    ("Shift", Inches(3.5), Inches(3.5), SLATE),
    ("Hyperscience", Inches(2.8), Inches(2.8), SLATE),
    ("FurtherAI", Inches(5.5), Inches(4.2), SLATE),
    ("Guidewire/\nDuck Creek", Inches(6.0), Inches(5.8), SLATE),
    ("CHARTERWELL", Inches(6.5), Inches(2.8), TEAL),
]
for name, cx, cy, color in competitors:
    # Dot
    dot_size = Inches(0.15) if name != "CHARTERWELL" else Inches(0.25)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, dot_size, dot_size)
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    # Label
    fs = 11 if name != "CHARTERWELL" else 14
    b = name == "CHARTERWELL"
    add_text_box(slide, cx + Inches(0.2), cy - Inches(0.1), Inches(1.5), Inches(0.5),
                 name, font_size=fs, font_color=color, bold=b)

# Right side summary
add_rect(slide, Inches(9.5), Inches(2.5), Inches(3.5), Inches(4.5), DEEP_NAVY)
add_multi_text(slide, Inches(9.8), Inches(2.7), Inches(3.0), Inches(4.0), [
    "Key differentiators:",
    "",
    "\u2022 Tractable: vision only, no docs",
    "\u2022 Shift: analytics DNA, no workspace",
    "\u2022 Hyperscience: horizontal IDP",
    "\u2022 FurtherAI: broad but shallow",
    "\u2022 Guidewire: legacy + AI bolt-on",
    "",
    "Charterwell = deep document",
    "intelligence + unified workspace.",
    "No one else occupies this quadrant.",
], font_size=13, font_color=WHITE)


# ============================================================
# SLIDE 9: MARKET
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DEEP_NAVY)

add_text_box(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.5),
             "MARKET OPPORTUNITY", font_size=14, font_color=TEAL, bold=True)
add_divider(slide, Inches(1), Inches(1.0), Inches(2), TEAL)

add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1),
             "Massive market. Massive white space.",
             font_size=36, font_color=WHITE, bold=True)

# Market size circles (concentric)
# TAM
add_text_box(slide, Inches(1), Inches(2.8), Inches(3.5), Inches(1.0),
             "$30B+", font_size=54, font_color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.0), Inches(3.5), Inches(0.5),
             "Claims operations market (TAM)", font_size=16, font_color=WHITE, alignment=PP_ALIGN.CENTER)

# SAM
add_text_box(slide, Inches(4.8), Inches(2.8), Inches(3.5), Inches(1.0),
             "$8\u201312B", font_size=54, font_color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(4.8), Inches(4.0), Inches(3.5), Inches(0.5),
             "Claims AI + automation (SAM)", font_size=16, font_color=WHITE, alignment=PP_ALIGN.CENTER)

# SOM
add_text_box(slide, Inches(8.8), Inches(2.8), Inches(3.5), Inches(1.0),
             "$500M\u2013$2B", font_size=42, font_color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(8.8), Inches(4.0), Inches(3.5), Inches(0.5),
             "Document-centric workspace (SOM)", font_size=16, font_color=WHITE, alignment=PP_ALIGN.CENTER)

# Key stats
add_rect(slide, Inches(1), Inches(5.0), Inches(11.3), Inches(2.0), RGBColor(0x14, 0x22, 0x3B))
add_multi_text(slide, Inches(1.5), Inches(5.2), Inches(10.3), Inches(1.7), [
    "\u2022  $570B in U.S. P&C premiums alone",
    "\u2022  88% of carriers want AI \u2014 only 7% have scaled it. The gap is the opportunity.",
    "\u2022  Insurance AI market growing 30%+ CAGR through 2030",
    "\u2022  Guidewire alone: $1.2B revenue, 19% ARR growth \u2014 proving carriers spend on claims technology",
    "\u2022  Category-defining vertical AI companies: Legora ($5.55B), Veeva ($50B+), Procore ($12B+)",
], font_size=15, font_color=WARM_WHITE)


# ============================================================
# SLIDE 10: BUSINESS MODEL
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WARM_WHITE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.5),
             "BUSINESS MODEL", font_size=14, font_color=TEAL, bold=True)
add_divider(slide, Inches(1), Inches(1.0), Inches(2))

add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
             "SaaS with built-in expansion.",
             font_size=36, font_color=DEEP_NAVY, bold=True)

# Pricing components
components = [
    ("Platform Fee", "Annual subscription for workspace\naccess, integrations, and\ncompliance module.\n\nEnsures predictable ARR."),
    ("Per-Claim Usage", "Usage-based fee aligned with\nclaims volume. Carriers already\nthink in per-claim economics.\n\nExpansion revenue built in."),
    ("Playbook Setup", "One-time professional services\nfor carrier-specific rule encoding.\n\nCreates switching-cost moat\nfrom institutional knowledge."),
]
for i, (title, desc) in enumerate(components):
    x = Inches(1 + i * 3.8)
    add_rect(slide, x, Inches(2.3), Inches(3.4), Inches(2.8), WHITE, SLATE)
    add_rect(slide, x, Inches(2.3), Inches(3.4), Pt(5), TEAL)
    add_text_box(slide, x + Inches(0.3), Inches(2.6), Inches(2.9), Inches(0.5),
                 title, font_size=20, font_color=DEEP_NAVY, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(3.3), Inches(2.9), Inches(1.8),
                 desc, font_size=14, font_color=CHARCOAL)

# ACV table
add_rect(slide, Inches(1), Inches(5.5), Inches(11.3), Inches(1.7), DEEP_NAVY)
add_multi_text(slide, Inches(1.5), Inches(5.6), Inches(10), Inches(1.5), [
    "Illustrative ACV by carrier size:",
    "",
    "Mid-market ($1\u20135B GWP):  $500K\u2013$1M    |    Large ($5\u201320B GWP):  $1M\u2013$3M    |    Top 20 ($20B+ GWP):  $3M\u2013$10M",
    "",
    "Land with a single claims line \u2192 expand across the carrier.",
], font_size=15, font_color=WHITE)


# ============================================================
# SLIDE 11: GO-TO-MARKET
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WARM_WHITE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.5),
             "GO-TO-MARKET", font_size=14, font_color=TEAL, bold=True)
add_divider(slide, Inches(1), Inches(1.0), Inches(2))

add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
             "Pilot-first. Prove it, then scale it.",
             font_size=36, font_color=DEEP_NAVY, bold=True)

# Timeline
phases = [
    ("Months 1\u20134", "DESIGN PARTNERS", "5 qualified conversations\n3 pilots launched\nGuide wire integration\nAnalyst briefing (Celent)"),
    ("Months 4\u20138", "PROOF OF VALUE", "60\u201390 day measurable pilots\nFirst case studies\nContent + thought leadership\nConference pipeline"),
    ("Months 8\u201312", "ENTERPRISE", "2\u20133 enterprise contracts\nVP Sales hired\nGuidewire Marketplace\nChannel partner conversations"),
    ("Year 2+", "SCALE", "Repeatable sales motion\nMulti-line expansion\nInternational\n$10M+ ARR target"),
]
for i, (time, title, desc) in enumerate(phases):
    x = Inches(0.6 + i * 3.2)
    # Phase card
    add_rect(slide, x, Inches(2.3), Inches(2.9), Inches(4.0), WHITE, SLATE)
    add_rect(slide, x, Inches(2.3), Inches(2.9), Pt(5), TEAL)
    add_text_box(slide, x + Inches(0.2), Inches(2.5), Inches(2.5), Inches(0.4),
                 time, font_size=13, font_color=TEAL, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(3.0), Inches(2.5), Inches(0.6),
                 title, font_size=18, font_color=DEEP_NAVY, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(3.7), Inches(2.5), Inches(2.3),
                 desc, font_size=13, font_color=CHARCOAL)

# Bottom
add_text_box(slide, Inches(1), Inches(6.6), Inches(11), Inches(0.5),
             "Target: mid-market and large P&C carriers  |  Guidewire/Duck Creek ecosystem  |  Founder-led sales in Year 1",
             font_size=14, font_color=SLATE, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 12: TEAM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WARM_WHITE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.5),
             "THE TEAM", font_size=14, font_color=TEAL, bold=True)
add_divider(slide, Inches(1), Inches(1.0), Inches(2))

add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
             "Insurance expertise meets AI expertise.",
             font_size=36, font_color=DEEP_NAVY, bold=True)

# Team placeholder cards
team_slots = [
    ("Founder & CEO", "[Name]", "Insurance industry + technology\nbackground. Deep understanding of\nclaims operations and the\ndocument processing pain."),
    ("Co-Founder & CTO", "[Name]", "AI/ML engineering expertise.\nDocument intelligence and NLP.\nScalable system architecture.\nEnterprise-grade infrastructure."),
    ("Advisors", "[Names]", "Insurance executives, AI researchers,\nenterprise SaaS operators.\nDomain credibility and\nnetwork access."),
]
for i, (role, name, desc) in enumerate(team_slots):
    x = Inches(1 + i * 3.8)
    add_rect(slide, x, Inches(2.5), Inches(3.4), Inches(3.5), WHITE, SLATE)
    add_rect(slide, x, Inches(2.5), Inches(3.4), Pt(5), TEAL)
    add_text_box(slide, x + Inches(0.3), Inches(2.8), Inches(2.9), Inches(0.4),
                 role, font_size=14, font_color=TEAL, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(3.3), Inches(2.9), Inches(0.5),
                 name, font_size=22, font_color=DEEP_NAVY, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(4.0), Inches(2.9), Inches(2.0),
                 desc, font_size=14, font_color=CHARCOAL)

# Why this team
add_rect(slide, Inches(1), Inches(6.3), Inches(11.3), Inches(0.9), DEEP_NAVY)
add_text_box(slide, Inches(1.5), Inches(6.4), Inches(10.3), Inches(0.7),
             "Why this team wins: the intersection of deep insurance domain knowledge and production-grade AI engineering. "
             "We\u2019ve lived the claims document problem and built the technology to solve it.",
             font_size=15, font_color=WHITE)


# ============================================================
# SLIDE 13: THE ASK
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DEEP_NAVY)

add_text_box(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.5),
             "THE ASK", font_size=14, font_color=TEAL, bold=True)
add_divider(slide, Inches(1), Inches(1.0), Inches(2), TEAL)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
             "Seed round: $[X]M",
             font_size=44, font_color=WHITE, bold=True)

add_text_box(slide, Inches(1), Inches(2.7), Inches(11), Inches(0.5),
             "What it funds:", font_size=22, font_color=TEAL, bold=True)

# Use of funds
funds = [
    ("Engineering", "Build the document intelligence engine\nand workspace product to production-readiness."),
    ("Design Partners", "Fund 3\u20135 carrier pilots, Guidewire integration,\nand first enterprise deployments."),
    ("Go-to-Market", "Hire founding sales + marketing.\nConference presence. Analyst engagement."),
    ("18-Month Runway", "Capital-efficient path to $1\u20133M ARR\nand Series A readiness."),
]
for i, (title, desc) in enumerate(funds):
    y = Inches(3.3 + i * 1.0)
    add_rect(slide, Inches(1.2), y, Inches(0.15), Inches(0.6), TEAL)
    add_text_box(slide, Inches(1.6), y, Inches(3), Inches(0.5),
                 title, font_size=18, font_color=WHITE, bold=True)
    add_text_box(slide, Inches(5), y, Inches(7), Inches(0.6),
                 desc, font_size=15, font_color=WARM_WHITE)


# ============================================================
# SLIDE 14: CLOSING
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DEEP_NAVY)

# Accent line
add_rect(slide, Inches(1.5), Inches(2.2), Inches(1.5), Pt(4), TEAL)

add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.2),
             "CHARTERWELL", font_size=54, font_color=WHITE, bold=True)

add_text_box(slide, Inches(1.5), Inches(3.7), Inches(10), Inches(0.8),
             "Every document understood.\nEvery decision informed.",
             font_size=28, font_color=TEAL, bold=False)

add_text_box(slide, Inches(1.5), Inches(5.2), Inches(6), Inches(0.4),
             "The AI-Native Claims Intelligence Workspace",
             font_size=18, font_color=WARM_WHITE)

# Contact
add_rect(slide, Inches(1.5), Inches(5.9), Inches(4), Inches(0.04), TEAL)
add_text_box(slide, Inches(1.5), Inches(6.1), Inches(6), Inches(0.8),
             "[Founder Name]  |  [email]  |  charterwell.ai",
             font_size=16, font_color=WARM_WHITE)

# Bottom bar
add_rect(slide, Inches(0), Inches(7.0), SLIDE_WIDTH, Inches(0.5), TEAL)
add_text_box(slide, Inches(1.5), Inches(7.05), Inches(6), Inches(0.4),
             "Confidential  |  March 2026", font_size=12, font_color=DEEP_NAVY, bold=False)


# ============================================================
# SAVE
# ============================================================
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, "charterwell-pitch-deck.pptx")
prs.save(output_path)
print(f"Deck saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
